#!/usr/bin/env python3
"""Create sample PDF, PPTX, and DOCX files for spike 004.

Generates test documents with embedded images, tables, headings,
and (for PPTX) speaker notes — covering the features that matter
most for conversion quality assessment.

Images are deliberately diverse to stress-test extraction quality:
  - Architecture diagram (boxes, arrows, labels — typical whiteboard diagram)
  - Bar chart with axes and labels (data-visualization style)
  - Photo-like gradient image with noise (simulates a photograph)

Usage:
    cd src/functions && uv run python ../spikes/004-pdf-pptx-conversion/create_samples.py
"""

from __future__ import annotations

import math
import random
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from reportlab.lib import colors
from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.lib.units import inch
from reportlab.platypus import (
    Image as RLImage,
    Paragraph,
    SimpleDocTemplate,
    Spacer,
    Table,
    TableStyle,
)

SAMPLES_DIR = Path(__file__).resolve().parent / "samples"


def _font() -> ImageFont.FreeTypeFont | ImageFont.ImageFont:
    """Return the best available font for labels."""
    for path in [
        "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
        "/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf",
        "/usr/share/fonts/TTF/DejaVuSans.ttf",
    ]:
        try:
            return ImageFont.truetype(path, 14)
        except OSError:
            continue
    return ImageFont.load_default()


def _font_small() -> ImageFont.FreeTypeFont | ImageFont.ImageFont:
    """Return a smaller font for secondary labels."""
    for path in [
        "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
        "/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf",
    ]:
        try:
            return ImageFont.truetype(path, 11)
        except OSError:
            continue
    return ImageFont.load_default()


def _create_architecture_diagram(filename: str, width: int = 600, height: int = 400) -> Path:
    """Create a realistic architecture diagram with boxes, arrows, and labels.

    Mimics a typical cloud-architecture diagram with named components,
    directional arrows, and color-coded service groups.
    """
    img = Image.new("RGB", (width, height), "#F5F5F5")
    draw = ImageDraw.Draw(img)
    font = _font()
    font_sm = _font_small()

    # Title
    draw.text((width // 2 - 100, 10), "System Architecture", fill="#1a1a1a", font=font)

    # --- Service boxes ---
    boxes = [
        (50, 70, 170, 130, "#4472C4", "Web App"),
        (220, 70, 370, 130, "#548235", "Agent API"),
        (420, 70, 560, 130, "#BF8F00", "AI Search"),
        (50, 190, 170, 250, "#C55A11", "Blob Storage"),
        (220, 190, 370, 250, "#7030A0", "OpenAI"),
        (420, 190, 560, 250, "#538DD5", "Cosmos DB"),
        (220, 310, 370, 370, "#A5A5A5", "Functions"),
    ]
    for x1, y1, x2, y2, color, label in boxes:
        draw.rounded_rectangle([x1, y1, x2, y2], radius=8, fill=color, outline="#333333", width=2)
        tw = draw.textlength(label, font=font_sm)
        tx = x1 + (x2 - x1 - tw) / 2
        ty = y1 + (y2 - y1 - 14) / 2
        draw.text((tx, ty), label, fill="white", font=font_sm)

    # --- Arrows (lines with arrowheads) ---
    arrows = [
        (170, 100, 220, 100),   # Web App → Agent API
        (370, 100, 420, 100),   # Agent API → AI Search
        (295, 130, 295, 190),   # Agent API → OpenAI
        (110, 130, 110, 190),   # Web App → Blob Storage
        (490, 130, 490, 190),   # AI Search → Cosmos DB
        (295, 250, 295, 310),   # OpenAI → Functions
        (170, 220, 220, 220),   # Blob Storage → OpenAI
    ]
    for x1, y1, x2, y2 in arrows:
        draw.line([(x1, y1), (x2, y2)], fill="#333333", width=2)
        # Simple arrowhead
        dx, dy = x2 - x1, y2 - y1
        length = math.sqrt(dx * dx + dy * dy)
        if length > 0:
            ux, uy = dx / length, dy / length
            px, py = -uy, ux  # perpendicular
            draw.polygon([
                (x2, y2),
                (x2 - 8 * ux + 4 * px, y2 - 8 * uy + 4 * py),
                (x2 - 8 * ux - 4 * px, y2 - 8 * uy - 4 * py),
            ], fill="#333333")

    # Legend
    draw.text((30, height - 30), "← data flow    ◼ Azure service", fill="#666", font=font_sm)

    path = SAMPLES_DIR / filename
    img.save(path)
    return path


def _create_bar_chart(filename: str, width: int = 500, height: int = 350) -> Path:
    """Create a bar chart with axes, labels, and a legend.

    Simulates a data-visualization image found in performance reports.
    """
    img = Image.new("RGB", (width, height), "white")
    draw = ImageDraw.Draw(img)
    font = _font()
    font_sm = _font_small()

    draw.text((width // 2 - 80, 8), "Conversion Latency (ms)", fill="#1a1a1a", font=font)

    # Chart area
    left, right, top, bottom = 70, width - 40, 50, height - 60
    draw.line([(left, top), (left, bottom)], fill="#333", width=2)
    draw.line([(left, bottom), (right, bottom)], fill="#333", width=2)

    # Y-axis ticks and labels
    max_val = 500
    for val in range(0, max_val + 1, 100):
        y = bottom - (val / max_val) * (bottom - top)
        draw.line([(left - 5, y), (left, y)], fill="#333", width=1)
        draw.text((left - 40, y - 7), str(val), fill="#333", font=font_sm)
        if val > 0:
            draw.line([(left, y), (right, y)], fill="#E0E0E0", width=1)

    # Data
    categories = ["CU", "Mistral", "MarkItDown", "MarkItDown\n+ PDF"]
    values_text = [320, 280, 45, 65]
    values_image = [180, 210, 180, 195]
    bar_colors_text = ["#4472C4", "#548235", "#BF8F00", "#C55A11"]
    bar_colors_img = ["#8FAADC", "#A9D18E", "#FFD966", "#F4B183"]

    n = len(categories)
    group_width = (right - left - 20) / n
    bar_width = group_width * 0.35

    for i, (cat, vt, vi) in enumerate(zip(categories, values_text, values_image)):
        gx = left + 10 + i * group_width
        # Text extraction bar
        bh_t = (vt / max_val) * (bottom - top)
        draw.rectangle([gx, bottom - bh_t, gx + bar_width, bottom], fill=bar_colors_text[i])
        # Image extraction bar
        bh_i = (vi / max_val) * (bottom - top)
        draw.rectangle(
            [gx + bar_width + 2, bottom - bh_i, gx + 2 * bar_width + 2, bottom],
            fill=bar_colors_img[i],
        )
        # Category label
        draw.text((gx, bottom + 5), cat, fill="#333", font=font_sm)

    # Legend
    lx = right - 160
    draw.rectangle([lx, top, lx + 12, top + 12], fill="#4472C4")
    draw.text((lx + 16, top - 1), "Text extraction", fill="#333", font=font_sm)
    draw.rectangle([lx, top + 18, lx + 12, top + 30], fill="#8FAADC")
    draw.text((lx + 16, top + 17), "Image extraction", fill="#333", font=font_sm)

    path = SAMPLES_DIR / filename
    img.save(path)
    return path


def _create_photo_like(filename: str, width: int = 500, height: int = 350) -> Path:
    """Create a photo-like image with gradient, noise, and shapes.

    Simulates a photograph to test extraction of non-diagram images —
    includes color gradients, random noise, and organic shapes.
    """
    img = Image.new("RGB", (width, height))
    pixels = img.load()
    rng = random.Random(42)  # reproducible

    # Gradient background (sky-to-ground)
    for y in range(height):
        ratio = y / height
        r = int(135 * (1 - ratio) + 34 * ratio)
        g = int(206 * (1 - ratio) + 139 * ratio)
        b = int(235 * (1 - ratio) + 34 * ratio)
        for x in range(width):
            nr = max(0, min(255, r + rng.randint(-8, 8)))
            ng = max(0, min(255, g + rng.randint(-8, 8)))
            nb = max(0, min(255, b + rng.randint(-8, 8)))
            pixels[x, y] = (nr, ng, nb)

    draw = ImageDraw.Draw(img)

    # "Trees" — filled green ellipses in the lower half
    for _ in range(6):
        cx = rng.randint(30, width - 30)
        cy = rng.randint(height // 2, height - 40)
        rw = rng.randint(30, 70)
        rh = rng.randint(40, 80)
        green = rng.randint(80, 160)
        draw.ellipse([cx - rw, cy - rh, cx + rw, cy + rh], fill=(30, green, 20))

    # "Sun" — yellow circle in upper-right
    draw.ellipse([width - 100, 20, width - 30, 90], fill=(255, 230, 80), outline=(255, 200, 50), width=3)

    # Some "clouds" — semi-transparent white ellipses
    for _ in range(3):
        cx = rng.randint(50, width - 100)
        cy = rng.randint(20, height // 4)
        draw.ellipse([cx, cy, cx + rng.randint(60, 120), cy + rng.randint(20, 40)], fill=(240, 240, 255))

    path = SAMPLES_DIR / filename
    img.save(path)
    return path


def create_sample_pdf() -> Path:
    """Create a sample PDF with headings, tables, images, and links.

    Uses three distinct image types to stress-test extraction:
      1. Architecture diagram (boxes + arrows + labels)
      2. Bar chart (axes + data bars + legend)
      3. Photo-like (gradient + noise + organic shapes)
    """
    SAMPLES_DIR.mkdir(parents=True, exist_ok=True)

    # Create diverse test images
    img_arch = _create_architecture_diagram("diagram-architecture.png", 600, 400)
    img_chart = _create_bar_chart("chart-performance.png", 500, 350)
    img_photo = _create_photo_like("photo-landscape.png", 500, 350)

    pdf_path = SAMPLES_DIR / "sample-article.pdf"
    doc = SimpleDocTemplate(str(pdf_path), pagesize=letter)
    styles = getSampleStyleSheet()
    story: list = []

    # Title
    story.append(Paragraph("Azure Knowledge Base Architecture Guide", styles["Title"]))
    story.append(Spacer(1, 0.25 * inch))

    # Introduction paragraph
    story.append(Paragraph("Introduction", styles["Heading1"]))
    story.append(Paragraph(
        "This document describes the architecture of the Azure Knowledge Base system. "
        "The system uses Azure AI Search for retrieval-augmented generation (RAG) and "
        "Azure OpenAI Service for natural language understanding.",
        styles["Normal"],
    ))
    story.append(Spacer(1, 0.15 * inch))

    # Embedded image 1
    story.append(Paragraph("System Architecture", styles["Heading2"]))
    story.append(Paragraph(
        "The following diagram shows the high-level architecture of the system:",
        styles["Normal"],
    ))
    story.append(Spacer(1, 0.1 * inch))
    story.append(RLImage(str(img_arch), width=5 * inch, height=3.3 * inch))
    story.append(Spacer(1, 0.15 * inch))

    # Table
    story.append(Paragraph("Component Overview", styles["Heading2"]))
    table_data = [
        ["Component", "Service", "Purpose", "SKU"],
        ["Search Index", "Azure AI Search", "Vector + keyword search", "Standard S1"],
        ["LLM", "Azure OpenAI", "GPT-4.1 for generation", "Standard"],
        ["Storage", "Azure Blob Storage", "KB article files", "Standard LRS"],
        ["Identity", "Azure Entra ID", "Authentication & RBAC", "P1"],
        ["Monitoring", "Application Insights", "Telemetry & logging", "Pay-as-you-go"],
    ]
    table = Table(table_data, colWidths=[1.2 * inch, 1.5 * inch, 2 * inch, 1.3 * inch])
    table.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#4472C4")),
        ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
        ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
        ("GRID", (0, 0), (-1, -1), 0.5, colors.grey),
        ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor("#D9E2F3")]),
    ]))
    story.append(table)
    story.append(Spacer(1, 0.15 * inch))

    # More content with sub-headings
    story.append(Paragraph("Deployment Options", styles["Heading2"]))
    story.append(Paragraph("Container Apps Hosting", styles["Heading3"]))
    story.append(Paragraph(
        "The agent is deployed as an Azure Container App with managed identity enabled. "
        "Environment variables are configured via Bicep templates, and secrets are stored "
        "in Azure Key Vault.",
        styles["Normal"],
    ))

    story.append(Paragraph("Scaling Configuration", styles["Heading3"]))
    story.append(Paragraph(
        "Auto-scaling is configured based on HTTP request concurrency. The minimum replica "
        "count is set to 1 for development environments and 3 for production.",
        styles["Normal"],
    ))

    # Second image
    story.append(Spacer(1, 0.15 * inch))
    story.append(Paragraph("Performance Metrics", styles["Heading2"]))
    story.append(Paragraph(
        "The chart below compares conversion latency across backends:",
        styles["Normal"],
    ))
    story.append(RLImage(str(img_chart), width=4.5 * inch, height=3.15 * inch))
    story.append(Spacer(1, 0.15 * inch))

    # Bullet points
    story.append(Paragraph("Key Features", styles["Heading2"]))
    bullets = [
        "Retrieval-augmented generation with hybrid search (vector + keyword)",
        "Image-grounded responses using GPT-4.1 vision capabilities",
        "Per-request context isolation using Python ContextVar",
        "Managed identity for all service-to-service authentication",
        "Infrastructure as Code using Bicep templates",
    ]
    for bullet in bullets:
        story.append(Paragraph(f"• {bullet}", styles["Normal"]))

    # Third image: photo-like (gradient + noise — simulates a photograph)
    story.append(Spacer(1, 0.15 * inch))
    story.append(Paragraph("Sample Photo", styles["Heading2"]))
    story.append(Paragraph(
        "Photo-like images with complex color gradients and organic shapes "
        "test whether image extraction preserves color fidelity and detail:",
        styles["Normal"],
    ))
    story.append(RLImage(str(img_photo), width=4.5 * inch, height=3.15 * inch))

    # Second table
    story.append(Spacer(1, 0.15 * inch))
    story.append(Paragraph("API Endpoints", styles["Heading2"]))
    api_table_data = [
        ["Endpoint", "Method", "Description"],
        ["/api/chat", "POST", "Send a chat message to the agent"],
        ["/api/health", "GET", "Health check endpoint"],
        ["/api/convert-markitdown", "POST", "Convert HTML to Markdown"],
        ["/api/index", "POST", "Index articles into AI Search"],
    ]
    api_table = Table(api_table_data, colWidths=[2 * inch, 1 * inch, 3 * inch])
    api_table.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#4472C4")),
        ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
        ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
        ("GRID", (0, 0), (-1, -1), 0.5, colors.grey),
    ]))
    story.append(api_table)

    # Build PDF
    doc.build(story)
    print(f"  Created PDF: {pdf_path} ({pdf_path.stat().st_size:,} bytes)")
    return pdf_path


def create_sample_pptx() -> Path:
    """Create a sample PPTX with slides, speaker notes, images, and tables."""
    from pptx import Presentation
    from pptx.util import Inches, Pt

    SAMPLES_DIR.mkdir(parents=True, exist_ok=True)

    # Create diverse test images
    img_arch = _create_architecture_diagram("slide-diagram.png", 600, 400)
    img_chart = _create_bar_chart("slide-chart.png", 500, 350)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title slide
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Azure KB Agent Architecture"
    slide1.placeholders[1].text = "PDF/PPTX Conversion Quality Spike\nMarch 2026"
    notes1 = slide1.notes_slide
    notes1.notes_text_frame.text = (
        "Welcome to the architecture overview for the Azure KB Agent. "
        "This presentation covers the key components, deployment strategy, "
        "and performance characteristics of the system. "
        "IMPORTANT: Mention that we use managed identity everywhere — no secrets in code."
    )

    # Slide 2: Content with bullet points
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "System Components"
    body2 = slide2.placeholders[1]
    tf2 = body2.text_frame
    tf2.text = "Core Services"
    for item in [
        "Azure AI Search — hybrid vector + keyword retrieval",
        "Azure OpenAI — GPT-4.1 for generation and vision",
        "Azure Blob Storage — KB article staging and serving",
        "Azure Container Apps — serverless hosting",
        "Azure Cosmos DB — conversation memory",
    ]:
        p = tf2.add_paragraph()
        p.text = item
        p.level = 1
    notes2 = slide2.notes_slide
    notes2.notes_text_frame.text = (
        "Emphasize that all services communicate via managed identity. "
        "No connection strings or API keys are stored in application code. "
        "The Container Apps use DefaultAzureCredential for authentication. "
        "Cosmos DB uses native RBAC with Built-in Data Contributor role."
    )

    # Slide 3: Image slide
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    slide3.shapes.title.text = "Architecture Diagram"
    slide3.shapes.add_picture(str(img_arch), Inches(2), Inches(2), Inches(9), Inches(4.5))
    notes3 = slide3.notes_slide
    notes3.notes_text_frame.text = (
        "This diagram shows the data flow from HTML ingestion through "
        "conversion (MarkItDown), indexing (AI Search), to the agent query path. "
        "Key point: the conversion pipeline supports three backends — "
        "Content Understanding, Mistral Document AI, and MarkItDown."
    )

    # Slide 4: Table slide
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])
    slide4.shapes.title.text = "Conversion Backend Comparison"
    rows, cols = 5, 4
    table_shape = slide4.shapes.add_table(rows, cols, Inches(1), Inches(2), Inches(11), Inches(4))
    tbl = table_shape.table
    headers = ["Backend", "API Calls", "Image Handling", "Cost"]
    data = [
        ["Content Understanding", "CU API", "CU Analyzer", "CU + GPT-4.1"],
        ["Mistral Document AI", "Mistral OCR", "GPT-4.1 Vision", "Mistral + GPT-4.1"],
        ["MarkItDown", "None (local)", "GPT-4.1 Vision", "GPT-4.1 only"],
        ["MarkItDown + PDF", "TBD", "TBD", "TBD"],
    ]
    for ci, h in enumerate(headers):
        tbl.cell(0, ci).text = h
    for ri, row_data in enumerate(data, 1):
        for ci, val in enumerate(row_data):
            tbl.cell(ri, ci).text = val
    notes4 = slide4.notes_slide
    notes4.notes_text_frame.text = (
        "The MarkItDown backend is the most cost-effective option. "
        "The PDF row is for this spike — we are evaluating whether "
        "MarkItDown can handle PDF input directly with acceptable quality."
    )

    # Slide 5: Second image + conclusion
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])
    slide5.shapes.title.text = "Performance Results"
    slide5.shapes.add_picture(str(img_chart), Inches(3), Inches(2), Inches(7), Inches(4))
    notes5 = slide5.notes_slide
    notes5.notes_text_frame.text = (
        "Performance chart shows that MarkItDown is 10x faster than cloud-based "
        "alternatives for text extraction. Image description time is the same "
        "across all backends since they all use GPT-4.1 vision."
    )

    pptx_path = SAMPLES_DIR / "sample-presentation.pptx"
    prs.save(str(pptx_path))
    print(f"  Created PPTX: {pptx_path} ({pptx_path.stat().st_size:,} bytes)")
    return pptx_path


def create_sample_docx() -> Path:
    """Create a sample DOCX with headings, tables, and an image."""
    from docx import Document
    from docx.shared import Inches

    SAMPLES_DIR.mkdir(parents=True, exist_ok=True)
    img_path = _create_architecture_diagram("docx-diagram.png", 600, 400)

    doc = Document()
    doc.add_heading("Azure KB Agent Setup Guide", level=0)
    doc.add_paragraph(
        "This guide covers the setup and configuration of the Azure Knowledge Base Agent "
        "for local development and Azure deployment."
    )

    doc.add_heading("Prerequisites", level=1)
    doc.add_paragraph("Python 3.11+", style="List Bullet")
    doc.add_paragraph("Azure CLI with active subscription", style="List Bullet")
    doc.add_paragraph("Azure Developer CLI (azd)", style="List Bullet")
    doc.add_paragraph("Docker Desktop", style="List Bullet")

    doc.add_heading("Architecture Overview", level=1)
    doc.add_paragraph("The following diagram shows the deployment architecture:")
    doc.add_picture(str(img_path), width=Inches(4.5))

    doc.add_heading("Environment Variables", level=1)
    table = doc.add_table(rows=5, cols=3, style="Table Grid")
    headers = ["Variable", "Required", "Description"]
    for ci, h in enumerate(headers):
        table.cell(0, ci).text = h
    env_vars = [
        ("AI_SERVICES_ENDPOINT", "Yes", "Azure AI Services endpoint URL"),
        ("SEARCH_ENDPOINT", "Yes", "Azure AI Search endpoint URL"),
        ("SERVING_BLOB_ENDPOINT", "Yes", "Blob storage endpoint for serving"),
        ("COSMOS_ENDPOINT", "No", "Cosmos DB endpoint for memory"),
    ]
    for ri, (var, req, desc) in enumerate(env_vars, 1):
        table.cell(ri, 0).text = var
        table.cell(ri, 1).text = req
        table.cell(ri, 2).text = desc

    doc.add_heading("Deployment Steps", level=1)
    doc.add_paragraph("Run azd up to provision and deploy all resources.", style="List Number")
    doc.add_paragraph("Configure environment with azd env get-values.", style="List Number")
    doc.add_paragraph("Run make test to verify the deployment.", style="List Number")

    docx_path = SAMPLES_DIR / "sample-document.docx"
    doc.save(str(docx_path))
    print(f"  Created DOCX: {docx_path} ({docx_path.stat().st_size:,} bytes)")
    return docx_path


def main() -> None:
    print("Creating sample documents for Spike 004...")
    create_sample_pdf()
    create_sample_pptx()
    create_sample_docx()
    print("Done — all samples in", SAMPLES_DIR)


if __name__ == "__main__":
    main()
